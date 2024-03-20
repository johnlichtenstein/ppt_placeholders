"""
Show what's happening to pptx layouts and placeholders
"""
import pptx
import pandas as pd
def listLayouts():
    lD = {i: x.name for i, x in enumerate(pptx.Presentation().slide_layouts)}
    return lD
def listPlaceholder(slide):
    pD = {i: x.name for i, x in enumerate(slide.placeholders)}
    return pD
def d2R(d, key="key", val="val", const=False, lname=""):
    tD = {}
    if const:
        tD[lname] = [const] * len(d)
    tD[key] = list(d.keys())
    tD[val] = [d[x] for x in d.keys()]
    tR = pd.DataFrame(tD)
    return tR

def allLayouts():
    lD = listLayouts()
    tL = []
    prs = pptx.Presentation()
    for k in lD.keys():
        slide_layout = prs.slide_layouts[k]
        slide = prs.slides.add_slide(slide_layout)
        placeholderD = listPlaceholder(slide)
        placeholderR = d2R(placeholderD, key="placeholder", val="description" \
            , const=lD[k], lname="layout")
        tL.append(placeholderR)
    return pd.concat(tL)


# Press the green button in the gutter to run the script.
if __name__ == '__main__':
    print(allLayouts().to_markdown(index=False))
    # layoutD = listLayouts()
    # print(layoutD)
    # print()
    # prs = pptx.Presentation()
    # slide_layout = prs.slide_layouts[0]
    # slide = prs.slides.add_slide(slide_layout)
    # placeholderD = listPlaceholder(slide)
    # print(d2R(placeholderD, key="placeholder", val="description" \
    #           , const="0", lname="layout").to_markdown(index=False))

